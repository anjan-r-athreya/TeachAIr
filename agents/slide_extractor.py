import os
import fitz  # PyMuPDF
from pptx import Presentation
from base_agent import BaseAgent

class SlideExtractor(BaseAgent):
    """
    Extracts structured slide content from PDF or PPTX files.
    For each slide, it provides:
      - slide number
      - text content
      - summarized image descriptions (optional)
    """

    def __init__(self, model="gemini-2.0-flash"):
        super().__init__(model=model)

    def _extract_text_and_images_from_pptx(self, slide):
        text_content = []
        image_placeholders = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_content.append(shape.text.strip())
            elif shape.shape_type == 13:  # PICTURE
                image_placeholders.append("[Image detected]")

        return "\n".join(text_content), image_placeholders

    def _extract_text_and_images_from_pdf(self, page, page_number):
        text_content = page.get_text().strip()
        image_placeholders = [
            f"[Image {img_index} detected on page {page_number}]"
            for img_index, img in enumerate(page.get_images(), start=1)
        ]
        return text_content, image_placeholders

    def _summarize_images(self, image_placeholders):
        
        if not image_placeholders:
            return []

        # Batch prompt
        prompt = "Summarize the following images briefly:\n" + "\n".join(image_placeholders)
        summary = self.send_prompt(prompt)
        return [summary]  # returning as a list for consistent structure

    def extract_from_pptx(self, file_path: str, summarize_images=True):
        prs = Presentation(file_path)
        slides_data = []

        for idx, slide in enumerate(prs.slides, start=1):
            text, images = self._extract_text_and_images_from_pptx(slide)
            summarized_images = self._summarize_images(images) if summarize_images else []
            slides_data.append({
                "slide_number": idx,
                "text": text,
                "images": summarized_images
            })

        return slides_data

    def extract_from_pdf(self, file_path: str, summarize_images=True):
        doc = fitz.open(file_path)
        slides_data = []

        for idx, page in enumerate(doc, start=1):
            text, images = self._extract_text_and_images_from_pdf(page, idx)
            summarized_images = self._summarize_images(images) if summarize_images else []
            slides_data.append({
                "slide_number": idx,
                "text": text,
                "images": summarized_images
            })

        return slides_data

    def run_task(self, file_path: str, summarize_images=True):
        _, ext = os.path.splitext(file_path)
        ext = ext.lower()
        if ext == ".pptx":
            return self.extract_from_pptx(file_path, summarize_images=summarize_images)
        elif ext == ".pdf":
            return self.extract_from_pdf(file_path, summarize_images=summarize_images)
        else:
            raise ValueError("Unsupported file type. Only .pptx or .pdf allowed.")
